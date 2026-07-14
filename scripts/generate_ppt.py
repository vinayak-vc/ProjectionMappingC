from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_slide(prs, title, content_lines, is_ascii=False):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    
    content_placeholder = slide.placeholders[1]
    text_frame = content_placeholder.text_frame
    text_frame.clear() # clear default paragraphs
    
    for i, line in enumerate(content_lines):
        p = text_frame.add_paragraph()
        p.text = line
        if is_ascii:
            p.font.name = 'Courier New'
            p.font.size = Pt(12)
            p.level = 0
            # Set exact line spacing for ASCII to look right
            p.line_spacing = 1.0
        else:
            if line.startswith('  -'):
                p.level = 2
                p.text = line.replace('  -', '').strip()
            elif line.startswith('-'):
                p.level = 1
                p.text = line.replace('-', '').strip()
            else:
                p.level = 0
            p.font.size = Pt(18)

def main():
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Projection Mapping SDK"
    subtitle.text = "Comprehensive Features & Technical Overview\nVersion 1.0.0"

    # Slide 1: Overview
    add_slide(prs, "Overview (Non-Technical)", [
        "What is the Projection Mapping SDK?",
        "- An engine-agnostic toolset for real-time projection mapping.",
        "- Maps digital content onto complex real-world physical objects.",
        "- Provides real-time warping, edge blending, and projector calibration.",
        "- Works out-of-the-box with Unity and Unreal Engine 5.",
        "- Allows multi-projector setups with seamless overlaps."
    ])

    # Slide 2: Architecture Overview
    add_slide(prs, "Architecture Overview (Technical)", [
        " +---------------------------------------+",
        " |      Unity / Unreal Applications      |",
        " +-------------------+-------------------+",
        " |  C# (Unity API)   | Blueprints (UE5)  |",
        " +-------------------+-------------------+",
        " |      C-API Layer (PImpl, ABI-Safe)    |",
        " +-------------------+-------------------+",
        " |           C++ Native SDK              |",
        " | [Math] [Geometry] [Warp] [Calibration]|",
        " +---------------------------------------+"
    ], is_ascii=True)

    # Slide 3: Core Module
    add_slide(prs, "1. Core Module", [
        "- Robust Error Handling: No C++ exceptions cross the ABI boundary.",
        "- Context Lifecycle: Root Context object manages state without global singletons.",
        "- Thread-safe Logging: Injectable logging callbacks to capture C++ logs in the Engine console.",
        "- HandleRegistry: Generation-checked 64-bit handles for memory-safe C-API access.",
        "- Config Store: Thread-safe, typed Key/Value configuration system."
    ])

    # Slide 4: Math Library
    add_slide(prs, "2. Math Library", [
        "- Zero-dependency custom Math engine (no GLM in public headers).",
        "- SIMD-ready layouts (16-byte cache alignment).",
        "- Core Types: Vector2, Vector3, Vector4, Matrix4, Quaternion, Transform.",
        "- Geometry primitives: Ray, Plane, BoundingBox.",
        "- Sub-millisecond matrix multiplication optimized for AVX2/NEON."
    ])

    # Slide 5: Geometry Engine
    add_slide(prs, "3. Geometry Engine", [
        "- High-performance Mesh data structures.",
        "- Spatial Acceleration: BVH and KDTree implementations for rapid raycasting.",
        "- Curves & Surfaces: Bezier Curves, Splines, and UV Mapping utilities.",
        "- Mesh Optimizer: Face welding and smooth normal generation.",
        "- Thread-Safe: Designed specifically for multi-threaded access without race conditions."
    ])

    # Slide 6: Deformation & Warp Engine
    add_slide(prs, "4. Warp Engine & Deformations", [
        "- Real-time vertex manipulations using DeformationFields.",
        "- Support for GridWarp (lattice points) and BezierPatch deformations.",
        "- WarpNode hierarchy for complex projector transformations.",
        "- Parallel Map-Reduce Mesh Normals: Fully locks-free multi-core normal recalculation.",
        "- Branchless Evaluation: std::clamp used for branchless math to allow vectorization."
    ])

    # Slide 7: Edge Blending
    add_slide(prs, "5. Blend Engine", [
        "- Resolves overlapping projector boundaries.",
        "- Edge Blending Math: Gamma correction and soft-edge gradient falloffs.",
        "- Luminance Compensation: Matches black-levels across different projectors.",
        "- Mask Generator: Dynamically calculates physical masks to block stray light."
    ])

    # Slide 8: OpenCV Calibration
    add_slide(prs, "6. Calibration (OpenCV Integration)", [
        "- Fully integrated OpenCV (hidden behind PImpl so no DLL ABI spillage occurs).",
        "- Camera Intrinsics & Extrinsics calibration workflows.",
        "- GrayCode Decoder: Flashes structured light patterns to automatically map physical pixels to projector pixels.",
        "- Direct VideoCapture: C++ handles camera capture to prevent Unity/Unreal main thread freezing.",
        "- 3D Triangulation: Converts 2D physical pixels into exact 3D models."
    ])

    # Slide 9: Multithreading & Performance
    add_slide(prs, "7. Performance Optimizations", [
        "- std::execution::par_unseq is utilized extensively across hot-paths.",
        "- Matrix crunching scales automatically across all available CPU cores.",
        "- Cache-friendly memory layouts prevent CPU pipeline stalling.",
        "- 100x100 Grid Warps execute in less than 1ms."
    ])

    # Slide 10: Unity Workflow
    add_slide(prs, "Unity Integration & Workflow", [
        "- Installation: Direct drag-and-drop via Unity Package Manager.",
        "- Interactive Setup Wizard: Guided UI (Tools > Projection Mapping > Setup Wizard).",
        "- NativeBindings.cs: Safe C# interop using IntPtr handles.",
        "- GrayCode Calibration Window directly inside the Unity Editor.",
        "- Zero configuration required: The native DLL is pre-injected into the package."
    ])

    # Slide 11: Unreal Engine Workflow
    add_slide(prs, "Unreal Engine 5 Integration & Workflow", [
        "- Installation: Drop into Plugins/ProjectionMapping.",
        "- Blueprints: Native Blueprint nodes (e.g., 'Create Warp Node', 'Apply Deformation Field').",
        "- C++ Modules: Simply add 'ProjectionMappingSDK' to PublicDependencyModuleNames.",
        "- Works with Static Meshes and Procedural Mesh Components.",
        "- Live C++ native execution directly inside the Unreal Viewport."
    ])
    
    # Slide 12: Release & Packaging
    add_slide(prs, "Automated Release Packaging", [
        "- build_release.py automatically bumps version and compiles MSVC Binaries.",
        "- Automated Python scripts bundle the SDK, Unity Package, and Unreal Plugin into a clean ZIP.",
        "- Seamless dependency injection ensures end-users don't have to compile C++ themselves.",
        "- Integrated into GitHub Actions for CI/CD across Windows, Linux, and macOS."
    ])

    prs.save('ProjectionMapping_SDK_Features.pptx')
    print("Saved presentation to ProjectionMapping_SDK_Features.pptx")

if __name__ == '__main__':
    main()
